"""Editable English PPTX of the figures (one display item per slide: composite + each panel)."""
import os

from pptx import Presentation
from pptx.util import Inches, Pt

from common import FIGS_DIR, ensure_dirs

SLIDES = [
    ("fig_main.png", "Figure 1. Two universal pathways under scrutiny",
     "(A) Country fertility-transition onsets span ~185 years; the pooled change point (95% CI "
     "1946-1969) is a summary of asynchronous transitions. (B) Aligning each country at its own "
     "onset collapses the between-country spread of lambda (3.4x variance reduction). (C) Under "
     "event alignment only the Phase I conserved quantity tightens; Phase II does not. (D) "
     "Within single countries, regions at nearly identical e0 scatter off the 'universal' curve."),
    ("panel_A_onset.png", "Panel A. Transition onsets span ~185 years",
     "Distribution of country-specific fertility-transition onset years (Gapminder, n=185). Red "
     "band: bootstrap 95% CI of the single pooled calendar change point."),
    ("panel_B_collapse.png", "Panel B. Event alignment collapses the spread",
     "Between-country SD of the crude birth rate around the alignment point, in calendar time "
     "(1950) versus event time (each country's onset)."),
    ("panel_C_cv.png", "Panel C. Asymmetric universality of the two pathways",
     "Coefficient of variation of the two conserved quantities within the paper's fitting "
     "windows, comparing the calendar split with event alignment."),
    ("panel_D_subnational.png", "Panel D. Sub-national heterogeneity",
     "Region-level crude birth rate versus life expectancy (2019) for five decentralised states, "
     "against the fitted national Phase II isocline."),
    ("fig_time_rescaling.png", "Figure 5. Time-axis rescaling does not create universality",
     "Left: crude birth rate versus years since each country's transition onset. Right: versus "
     "duration-normalised time s=(t-onset)/D. Rescaling the time axis does not shrink the raw "
     "between-country spread of lambda beyond onset alignment; only the amplitude-normalised "
     "shape collapses, and the relationship at matched e0 is invariant to any monotone time warp."),
]


def main():
    ensure_dirs()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for img, title, caption in SLIDES:
        path = os.path.join(FIGS_DIR, img)
        if not os.path.exists(path):
            continue
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22); p.font.bold = True
        # image scaled to fit
        from PIL import Image
        w, h = Image.open(path).size
        maxw, maxh = Inches(12.0), Inches(5.2)
        ratio = min(maxw / w, maxh / h)
        iw, ih = int(w * ratio), int(h * ratio)
        left = int((prs.slide_width - iw) / 2)
        s.shapes.add_picture(path, left, Inches(1.05), width=iw, height=ih)
        cap = s.shapes.add_textbox(Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.95))
        cp = cap.text_frame; cp.word_wrap = True
        cp.paragraphs[0].text = caption
        cp.paragraphs[0].font.size = Pt(11)
    out = os.path.join(FIGS_DIR, "figures_EN.pptx")
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    main()
